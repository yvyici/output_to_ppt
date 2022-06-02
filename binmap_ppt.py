#binmap-excel|-pic|-ppt|
import re
from os import walk
import os
import shutil
import glob
import time

path_cwd = os.path.dirname(os.path.realpath(__file__))
#print('file_path===')
#print (file_path)
#print(ZZZ)
#path_cwd="C:/Users/marvi/test2"
#pathex = os.path.join(path_cwd, "CP1_yield_RCRGX8000_#1~#5.xlsx")
path_data='data2'
##path_data='G:\.shortcut-targets-by-id\1-K20VmkbFWnulFWAKr_3uOzYjwQLIHcO\Engineer- IoTMemory\03-Product and Test Engineer\00-Customer_CP\D004-ADUC\AD18F020A0\007_RCRGX8000_25\CP1\01_binmap_data'
#path_data='G:/.shortcut-targets-by-id/1-K20VmkbFWnulFWAKr_3uOzYjwQLIHcO/Engineer- IoTMemory/03-Product and Test Engineer/00-Customer_CP/D004-ADUC/AD18F020A0/007_RCRGX8000_25\CP1/01_binmap_data'

path_data=path_data.replace('\\','/')
print(path_data)

pathex = os.path.realpath(os.path.join(path_data, "CP1_format.xlsx"))
pathmap1 = os.path.realpath(os.path.join(path_data, "Map1.pptx"))

if os.path.isfile(pathex):
    print(pathex)
    os.unlink(pathex)

if os.path.isfile(pathmap1):
    print(pathmap1)
    os.unlink(pathmap1)

#print(ZZZ)
shutil.copyfile("CP1_format_ini_2.xlsx", pathex) 

shutil.copyfile("Map1_ini.pptx", pathmap1) 

png_path='png_data' 

# write the TXT file to excel(in the same sheet for the format)
def binmap(self, sheetname):
    import openpyxl
    path = os.path.join(path_cwd,'%s' %self)
    print('path==='+path)
    fopen = open(path, 'r',encoding='utf-8')
    lines = fopen.readlines()
    if os.path.isfile(pathex):
        print(pathex)
        os.unlink(pathex) 
        shutil.copyfile("CP1_format_ini_2.xlsx", pathex)    
    file = openpyxl.load_workbook(pathex)
    sheet = file['org']                               # assigne sheet 5
    i = 0
    for line in lines:
        line = line.strip('\n')
        line = line.replace("\t",",")
        line = line.split(',')
        
        for index in range(len(line)):
            s= line[index]
            try:
                s = float(s)
            except ValueError:
                pass
            sheet.cell(i+1, index+1, s)
        i = i + 1
    file.copy_worksheet(sheet)                  # if  you don't want to copy the sheet mark this three line
    sheetcopy = file['org Copy']                  #
    sheetcopy.title = '%s'%sheetname            #
    file.save(pathex)

# to search the path (use re to find the file name)
def findTXT():
    a = []
    mypath = path_cwd
#  for root, dirs, files in walk(mypath):
#    print()
#    for f in files:
#        if re.findall(r'(\S+\-\d{2}_result.TXT)', f):
#            a.append(f)
    name=os.path.join(path_data,'*_result.TXT')
    print(name)
    a=sorted(glob.glob(name))
#    print('a===')
#    print(a)
#    print(ZZZ)
    return a

# create new folder to store pic
def reset_folder():
    import os
    #fn = (re.findall(r'\S+\-',findTXT()[0])[0])
    #fn='png_data'
    path = os.path.join(path_data,"%s"%png_path)
   
    if os.path.exists(path):
        #os.makedirs(path)

        shutil.rmtree(path)
        #time.sleep(2)
    print('path==='+path)
    if os.path.exists(path):
        print('The file exists! '+path)
    #os.mkdir(path)    
    print(os.path.realpath(path))
    os.mkdir((os.path.realpath(path)))
    return png_path

#def nfolder():
#    import os
#    #fn = (re.findall(r'\S+\-',findTXT()[0])[0])
#    fn='png_data'    
#    return fn

# from excel export to pic

def file2waferno(files):
    print(files)
    waferno=[]
    for file in files:
        m=re.search(r'\-(\d+)_result.TXT',file)
        #print('m===')
        #print(m)
        waferno.append(m.group(1))
    return waferno


def extopic():
    import excel2img
    b=0
    print('findTXT()===')
    print(findTXT())
    waferno=file2waferno(findTXT())
    print(waferno)
    #print(ZZZ)
    while True:
        if b<len(findTXT()):
            binmap(findTXT()[b], b+1)

            png_file=os.path.join(path_data,"%s/pic%02d.png"%(png_path,int(waferno[b])))
            print('pngfile==='+png_file)
            excel2img.export_img(pathex,png_file, 'org', "A1:GF134")
            b = b+1
        else:
            break
    

# put the pic in ppt
def toppt():
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation(pathmap1)
    blank_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(blank_slide_layout)
    a = []
    mypath = os.path.join(path_data,'%s'%png_path)
    for root, dirs, files in walk(mypath):
        for f in files:
            if re.findall(r'(pic\d+)', f):
                a.append(f)
    d = 1
    i = 1
    j = 1
    x = 0.5
    y = 1.2
    b = 0
    while True:
        left = Inches(x)
        left2 = Inches(x-0.2)
        top = Inches(y)
        if b == len(a):
            break;
        elif i <5 and j<4:
            fa = a[b]
            img_path = os.path.join(path_data,'%s/%s'%(png_path,fa))
            pic = slide.shapes.add_picture(img_path, left, top, width=Inches(2.1), height=Inches(2.1))
            txBox = slide.shapes.add_textbox(left2, top, width=Inches(0.5), height=Inches(0.5))
            tf = txBox.text_frame
            tf.text = '#' + (re.findall(r'\d{2}', fa)[0])
            x = x+2.3
            d = d+1
            b= b+1
        elif i == 5:
            i= 0
            j= j+1
            y= y+2.1
            x= 0.5
        elif j>3:
            slide = prs.slides.add_slide(blank_slide_layout)
            i= 0
            j= 1
            x= 0.5
            y= 1.2
        i = i+1

    prs.save(pathmap1)


if __name__ == '__main__':
    tic = time.time()
    reset_folder()  
    extopic()       # from binmap export to photo
    toppt()         # add the picture to ppt
    toc = time.time()
    print(toc-tic)